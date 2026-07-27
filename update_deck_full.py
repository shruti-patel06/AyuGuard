import sys
import pptx
from pptx.util import Pt
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

prs = pptx.Presentation('AyuGuard_Prototype_Submission_Deck.pptx')

NAVY = RGBColor(20, 45, 85)
TEAL = RGBColor(13, 148, 136)
DARK = RGBColor(30, 41, 59)
GRAY = RGBColor(71, 85, 105)

# SLIDE 1
s1 = prs.slides[0]
for shape in s1.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = 'Participant Details'
        p0.font.bold = True
        p0.font.size = Pt(22)
        p0.font.color.rgb = NAVY
        
        p1 = tf.add_paragraph()
        p1.text = 'Participant Name: Shruti Patel'
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = DARK
        
        p2 = tf.add_paragraph()
        p2.text = 'Hackathon Track: AI for Better Living and Smarter Communities (Healthcare Access & Community Wellness)'
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        
        p3 = tf.add_paragraph()
        p3.text = 'Project Title: AyuGuard — Ambient Multi-Agent Caregiver Platform for Independent Elderly Living'
        p3.font.size = Pt(16)
        p3.font.bold = True
        p3.font.color.rgb = TEAL

# SLIDE 2
s2 = prs.slides[1]
for shape in s2.shapes:
    if shape.has_text_frame and 'Brief' in shape.text_frame.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = 'Brief About The Idea'
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = NAVY
    elif shape.has_text_frame and 'Challenge' in shape.text_frame.text:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'The Core Healthcare Challenge for Seniors'
        p0.font.bold = True
        p0.font.size = Pt(16)
        p0.font.color.rgb = TEAL
        
        p1 = tf.add_paragraph()
        p1.text = '• Independent Living Risks: Millions of seniors live alone. Subtle health declines (dehydration, early metabolic shifts, mild infection) develop gradually and are frequently ignored until late-stage emergency hospitalization.'
        p1.font.size = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = '• The Caregiving Gap: Family caregivers lack continuous visibility. Passive wearables capture simple raw vitals but miss qualitative symptoms, while home cameras are intrusive and violate elderly dignity.'
        p2.font.size = Pt(12)
        
        p3 = tf.add_paragraph()
        p3.text = "AyuGuard's Ambient Multi-Agent Solution"
        p3.font.bold = True
        p3.font.size = Pt(16)
        p3.font.color.rgb = TEAL
        
        p4 = tf.add_paragraph()
        p4.text = '• Conversational Care: Seniors or caregivers simply chat in plain language (English, Hindi, Hinglish). Powered by Google ADK & Vertex AI, AyuGuard extracts structured symptoms from casual stories.'
        p4.font.size = Pt(12)
        
        p5 = tf.add_paragraph()
        p5.text = '• Multi-Condition Care Reconciliation: When chronic conditions (Pre-diabetes) conflict with acute symptoms (Diarrhoea), specialized AI subagents collaborate to adapt meal plans safely and manage recovery de-escalation back to baseline.'
        p5.font.size = Pt(12)

# SLIDE 3
s3 = prs.slides[2]
for shape in s3.shapes:
    if shape.has_text_frame and 'How We' in shape.text_frame.text:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'How We Approached the Problem & Real-World Impact'
        p0.font.bold = True
        p0.font.size = Pt(20)
        p0.font.color.rgb = NAVY
        
        p1 = tf.add_paragraph()
        p1.text = '• Google ADK Multi-Agent Framework: Built using Google ADK 2.x and Vertex AI (gemini-2.5-flash-lite). Features specialized subagents: extraction, clinical dataset retrieval, and dietary reconciliation.'
        p1.font.size = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = '• Deterministic Safety Gate: Integrates a 14-day rolling trend window and Jaccard similarity scoring against a 4,921-row clinical dataset. The LLM never decides if a pattern is dangerous — deterministic math does.'
        p2.font.size = Pt(12)
        
        p3 = tf.add_paragraph()
        p3.text = '• Multimodal Data Vault: Ingests unstructured chat, scanned PDF lab reports, and prescriptions into Google Cloud Storage with Gemini AI parameter extraction.'
        p3.font.size = Pt(12)
        
        p4 = tf.add_paragraph()
        p4.text = 'Real-World Practical Impact'
        p4.font.bold = True
        p4.font.size = Pt(16)
        p4.font.color.rgb = TEAL
        
        p5 = tf.add_paragraph()
        p5.text = '• Proactive Early Triage: Evaluates symptom frequency, severity, and multi-week trends to output a clear urgency score (Low, Watch, Escalate), preventing emergency room crises.'
        p5.font.size = Pt(12)

# SLIDE 4
s4 = prs.slides[3]
for shape in s4.shapes:
    if shape.has_text_frame and ('Differentiation' in shape.text_frame.text or 'Passive' in shape.text_frame.text):
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'Differentiation & Unique Selling Proposition (USP)'
        p0.font.bold = True
        p0.font.size = Pt(20)
        p0.font.color.rgb = NAVY
        
        p1 = tf.add_paragraph()
        p1.text = '• Passive Wearables -> Track vitals but miss qualitative symptoms; AyuGuard understands natural language.'
        p1.font.size = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = '• Intrusive Home Cameras -> Cause stress & privacy concerns; AyuGuard is 100% ambient, conversational, and dignity-preserving.'
        p2.font.size = Pt(12)
        
        p3 = tf.add_paragraph()
        p3.text = '• Standard Medical Chatbots -> Give generic replies without clinical grounding; AyuGuard tracks 14-day trends & clinical tables.'
        p3.font.size = Pt(12)
        
        p4 = tf.add_paragraph()
        p4.text = 'Unique Selling Proposition (USP)'
        p4.font.bold = True
        p4.font.size = Pt(16)
        p4.font.color.rgb = TEAL
        
        p5 = tf.add_paragraph()
        p5.text = '• Multi-Condition Subagent Collaboration: Resolves dietary conflicts (Pre-diabetes high fiber vs Diarrhoea low fiber) and handles recovery de-escalation back to baseline.'
        p5.font.size = Pt(12)
        
        p6 = tf.add_paragraph()
        p6.text = '• Real-Time Synchronized Engine: Glassmorphism SPA with background polling (3s notifications, 5s dashboard sync) connecting Caregiver & Patient views.'
        p6.font.size = Pt(12)

# SLIDE 5
s5 = prs.slides[4]
for shape in s5.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'List of Features Offered by AyuGuard'
        p0.font.bold = True
        p0.font.size = Pt(20)
        p0.font.color.rgb = NAVY
        
        feats = [
            '1. Ambient Conversational Logging: Multilingual NLP parser for English, Hindi, and Hinglish.',
            '2. Multi-Condition Dietary Reconciliation Subagent: Resolves conflicting health rules (Pre-diabetes + Diarrhoea) & handles recovery lifecycle transitions.',
            '3. 14-Day Deterministic Trend Scorer: Calculates composite urgency (Low, Watch, Escalate) using clinical dataset weighting (4,921 rows).',
            '4. Multimodal Medical Record Vault: Gemini AI parses PDF/image lab reports, extracts abnormal values (HbA1c, glucose), and tracks parameters.',
            '5. Real-Time Dual Persona SPA: Glassmorphism UI with real-time background sync between Caregiver Dashboard and Patient Chat.',
            '6. Instant History Reset & Role Context Isolation: Instant chat wipe + context guard ensuring caregiver is addressed as Priya and patient as Rajan ji.'
        ]
        for f in feats:
            p = tf.add_paragraph()
            p.text = f
            p.font.size = Pt(12)

# SLIDE 6
s6 = prs.slides[5]
for shape in s6.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'Process Flow & Multi-Agent Pipeline Architecture'
        p0.font.bold = True
        p0.font.size = Pt(20)
        p0.font.color.rgb = NAVY
        
        flow_text = [
            '1. Input Ingestion -> Caregiver or Patient chats in English, Hindi, or Hinglish via FastAPI SSE Proxy.',
            '2. Orchestrator Dispatch -> ayuguard_orchestrator (Root ADK Agent) evaluates active persona context.',
            '3. Symptom Extraction -> symptom_extraction_agent converts free text into structured symptom JSON.',
            '4. Deterministic Urgency Scorer -> 14-day rolling trend window calculates composite urgency score (Low, Watch, Escalate).',
            '5. Subagent Reconciliation -> dietary_reconciliation_agent resolves multi-condition conflicts (Pre-diabetes + Diarrhoea) & recovery transitions.',
            '6. Real-Time UI Dispatch -> save_care_plan() updates local JSON/Firestore & triggers real-time polling updates on the SPA Dashboard.'
        ]
        for f in flow_text:
            p = tf.add_paragraph()
            p.text = f
            p.font.size = Pt(12)

# SLIDE 7
s7 = prs.slides[6]
for shape in s7.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'Responsive Single-Page Application (SPA) Layout'
        p0.font.bold = True
        p0.font.size = Pt(20)
        p0.font.color.rgb = NAVY
        
        layout_text = [
            '• Caregiver Chat Tab: Natural language symptom logging, AI subagent advice, and instant history clear button.',
            '• Health Dashboard Tab: Glassmorphism UI featuring 14-day Urgency Ring, Symptom Timeline, Frequency Bars, Disease Match, and Active Care Plan.',
            '• Patient Chat View: Calming direct patient interface with real-time in-app notification cards.',
            '• Medical Records Vault: Drag-and-drop PDF/image upload zone with Gemini multimodal parameter extraction.'
        ]
        for l in layout_text:
            p = tf.add_paragraph()
            p.text = l
            p.font.size = Pt(12)

# SLIDE 8
s8 = prs.slides[7]
for shape in s8.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'Live Production Deployments & Code Repository'
        p0.font.bold = True
        p0.font.size = Pt(20)
        p0.font.color.rgb = NAVY
        
        urls = [
            '• Live SPA Dashboard: https://ayuguard-ui-118454190848.asia-south1.run.app',
            '• Live Agent Orchestrator API: https://ayuguard-agent-118454190848.asia-south1.run.app',
            '• GitHub Repository: https://github.com/shruti-patel06/AyuGuard',
            '• Cloud Infrastructure: Google Cloud Run (asia-south1) | GCP Project: silken-dogfish-484814-g9'
        ]
        for u in urls:
            p = tf.add_paragraph()
            p.text = u
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = TEAL

# SLIDE 9
s9 = prs.slides[8]
for shape in s9.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'Technologies Used & System Design Justification'
        p0.font.bold = True
        p0.font.size = Pt(20)
        p0.font.color.rgb = NAVY
        
        techs = [
            '• Google ADK (Agent Development Kit): Modular Python framework establishing strict subagent boundaries and reliable tool execution loops.',
            '• Vertex AI (Gemini 2.5 Flash Lite): Enterprise LLM powering NLP symptom extraction, dietary reconciliation reasoning, and multimodal document analysis in us-central1.',
            '• Google Cloud Run: Managed serverless container hosting providing auto-scaling from zero to high concurrency in asia-south1.',
            '• Google Cloud Storage & Firestore: Encrypted PDF document vault and real-time notification feed sync.',
            '• Google Secret Manager: Secure API key & ADC credential management with enterprise compliance.'
        ]
        for t in techs:
            p = tf.add_paragraph()
            p.text = t
            p.font.size = Pt(12)

# SLIDE 10
s10 = prs.slides[9]
for shape in s10.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'Snapshots of Prototype: Caregiver Chat & Subagent Collaboration'
        p0.font.bold = True
        p0.font.size = Pt(18)
        p0.font.color.rgb = NAVY
        
        snaps1 = [
            '• Conversational Logging: Accepts English, Hindi, and Hinglish input (e.g. "Papa ko loose motion ho gaya hai").',
            '• Subagent Collaboration: dietary_reconciliation_agent automatically detects dietary conflicts and updates the care plan.',
            '• Real-Time Notification Bell: Flashes instantly when care plan revisions or critical alerts occur.'
        ]
        for s in snaps1:
            p = tf.add_paragraph()
            p.text = s
            p.font.size = Pt(12)

# SLIDE 11
s11 = prs.slides[10]
for shape in s11.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'Snapshots of Prototype: Health Dashboard & Medical Records Vault'
        p0.font.bold = True
        p0.font.size = Pt(18)
        p0.font.color.rgb = NAVY
        
        snaps2 = [
            '• Glassmorphism Health Dashboard: Visual composite Urgency Ring, 14-day chronological timeline, and clinical precautions.',
            '• Multimodal Gemini Records Vault: Automatic extraction of abnormal lab parameters (HbA1c, blood glucose) from uploaded PDF reports.',
            '• Patient Chat View: Calm patient interface with real-time caregiver update cards.'
        ]
        for s in snaps2:
            p = tf.add_paragraph()
            p.text = s
            p.font.size = Pt(12)

# SLIDE 12
s12 = prs.slides[11]
for shape in s12.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = 'Conclusion & Future Roadmap'
        p0.font.bold = True
        p0.font.size = Pt(22)
        p0.font.color.rgb = NAVY
        
        conc = [
            '• Clinical Safety Boundary: AyuGuard supplements — never replaces — professional medical advice ("The math decides what is dangerous — the LLM only explains it").',
            '• Future Roadmap: Integration with smart IoT pill dispensers, voice-call IVR integration for non-smartphone seniors, and multi-caregiver family sync.',
            '• AyuGuard — Ambient Care, Powered by AI.'
        ]
        for c in conc:
            p = tf.add_paragraph()
            p.text = c
            p.font.size = Pt(14)

prs.save('AyuGuard_Prototype_Submission_Deck.pptx')
print('FULL 12 SLIDES UPDATE SUCCESSFUL!')

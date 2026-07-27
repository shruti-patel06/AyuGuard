import sys
import os
import pptx
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

prs = pptx.Presentation('AyuGuard_Prototype_Submission_Deck.pptx')

NAVY = RGBColor(20, 45, 85)
TEAL = RGBColor(13, 148, 136)
DARK = RGBColor(30, 41, 59)
GRAY = RGBColor(71, 85, 105)

def update_slide(slide, title, bullets, header_color=NAVY):
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    for s in text_shapes:
        s.text_frame.clear()
    
    if not text_shapes:
        target_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8.5), Inches(5.5))
    else:
        target_shape = text_shapes[-1]
    
    tf = target_shape.text_frame
    
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.bold = True
    p0.font.size = Pt(20)
    p0.font.color.rgb = header_color
    
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item[0] if isinstance(item, tuple) else item
        p.font.size = Pt(item[1] if isinstance(item, tuple) else 12)
        if isinstance(item, tuple) and len(item) > 2 and item[2]:
            p.font.bold = True
        if isinstance(item, tuple) and len(item) > 3:
            p.font.color.rgb = item[3]

# SLIDE 1
update_slide(
    prs.slides[0],
    "Participant Details",
    [
        ("Participant Name: Shruti Patel", 15, True, DARK),
        ("Hackathon Track: AI for Better Living and Smarter Communities (Healthcare Access & Community Wellness)", 14, False, GRAY),
        ("Project Title: AyuGuard — Ambient Multi-Agent Caregiver Platform for Independent Elderly Living", 16, True, TEAL)
    ]
)

# SLIDE 2
update_slide(
    prs.slides[1],
    "Brief About The Idea",
    [
        ("The Core Healthcare Challenge for Seniors", 15, True, TEAL),
        "• Independent Living Risks: Millions of seniors live alone. Subtle health declines (dehydration, early metabolic shifts, mild infection) develop gradually and are frequently ignored until late-stage emergency hospitalization.",
        "• The Caregiving Gap: Family caregivers lack continuous visibility. Passive wearables capture simple raw vitals but miss qualitative symptoms, while home cameras are intrusive and violate elderly dignity.",
        ("AyuGuard's Ambient Multi-Agent Solution", 15, True, TEAL),
        "• Conversational Care: Seniors or caregivers simply chat in plain language (English, Hindi, Hinglish). Powered by Google ADK & Vertex AI, AyuGuard extracts structured symptoms from casual stories.",
        "• Multi-Condition Care Reconciliation: When chronic conditions (Pre-diabetes) conflict with acute symptoms (Diarrhoea), specialized AI subagents collaborate to adapt meal plans safely and manage recovery de-escalation back to baseline."
    ]
)

# SLIDE 3
s3 = prs.slides[2]
update_slide(
    s3,
    "How We Approached the Problem & Decision Architecture",
    [
        ("Google ADK & Decision Logic Architecture", 13, True, TEAL),
        "• Deterministic Urgency Gate: 14-day rolling decay window + similarity search against 4,921 clinical rows.",
        "• Triage Decision Routing: Score >= 0.65 -> Escalate (Doctor), Score >= 0.42 -> Watch, Score < 0.42 -> Low.",
        "• Subagent Conflict Gate: Pre-diabetes + Diarrhoea -> Reconcile to temporary bland diet + ORS.",
        "• Recovery Transition: Symptoms resolved -> Restore healthy baseline meal plan."
    ]
)
img3_path = r"C:\Users\Shruti\.gemini\antigravity-ide\brain\4e55b1e4-91df-421f-8cc1-0a0d55313f34\ayuguard_decision_workflow_diagram_1785177163194.png"
if os.path.exists(img3_path):
    s3.shapes.add_picture(img3_path, Inches(5.0), Inches(1.6), width=Inches(4.5))

# SLIDE 4
update_slide(
    prs.slides[3],
    "Differentiation & Unique Selling Proposition (USP)",
    [
        ("Key Differentiators", 15, True, TEAL),
        "• Passive Wearables -> Track vitals but miss qualitative symptoms; AyuGuard understands natural language.",
        "• Intrusive Home Cameras -> Cause stress & privacy concerns; AyuGuard is 100% ambient, conversational, and dignity-preserving.",
        "• Standard Medical Chatbots -> Give generic replies without clinical grounding; AyuGuard tracks 14-day trends & clinical tables.",
        ("Unique Selling Proposition (USP)", 15, True, TEAL),
        "• Multi-Condition Subagent Collaboration: Resolves dietary conflicts (Pre-diabetes high fiber vs Diarrhoea low fiber) and handles recovery de-escalation back to baseline.",
        "• Real-Time Synchronized Engine: Glassmorphism SPA with background polling (3s notifications, 5s dashboard sync) connecting Caregiver & Patient views."
    ]
)

# SLIDE 5
update_slide(
    prs.slides[4],
    "List of Features Offered by AyuGuard",
    [
        "1. Ambient Conversational Logging: Multilingual NLP parser for English, Hindi, and Hinglish.",
        "2. Multi-Condition Dietary Reconciliation Subagent: Resolves conflicting health rules (Pre-diabetes + Diarrhoea) & handles recovery lifecycle transitions.",
        "3. 14-Day Deterministic Trend Scorer: Calculates composite urgency (Low, Watch, Escalate) using clinical dataset weighting (4,921 rows).",
        "4. Multimodal Medical Record Vault: Gemini AI parses PDF/image lab reports, extracts abnormal values (HbA1c, glucose), and tracks parameters.",
        "5. Real-Time Dual Persona SPA: Glassmorphism UI with real-time background sync between Caregiver Dashboard and Patient Chat.",
        "6. Instant History Reset & Role Context Isolation: Instant chat wipe + context guard ensuring caregiver is addressed as Priya and patient as Rajan ji."
    ]
)

import os

# SLIDE 6 - CASE Tools & System Use-Case Architecture
s6 = prs.slides[5]
update_slide(
    s6,
    "CASE Tools & System Use-Case Architecture Diagram",
    [
        ("Actors: Caregiver (Priya), Senior Patient (Rajan ji), ADK AI System Engine", 13, True, TEAL),
        "• UC-1: Multilingual Symptom Ingestion -> Free text to JSON (symptom, severity, date).",
        "• UC-2: Longitudinal Urgency Triage -> 14-day rolling window + 4,921-row dataset search.",
        "• UC-3: Multi-Condition Dietary Reconciliation -> Resolves Pre-diabetes + Diarrhoea conflicts.",
        "• UC-4: Multimodal Record Vault -> Gemini Vision lab report parameter extraction."
    ]
)
img1_path = r"C:\Users\Shruti\.gemini\antigravity-ide\brain\4e55b1e4-91df-421f-8cc1-0a0d55313f34\ayuguard_case_use_case_diagram_1785176915006.png"
if os.path.exists(img1_path):
    s6.shapes.add_picture(img1_path, Inches(5.0), Inches(1.6), width=Inches(4.5))

# SLIDE 7 - User Workflow Diagram
s7 = prs.slides[6]
update_slide(
    s7,
    "End-to-End User Workflow & Journey Diagram",
    [
        ("Phase 1: Input Stage", 13, True, TEAL),
        "• Caregiver logs symptoms or uploads lab PDF in natural language.",
        ("Phase 2: ADK Multi-Agent Stage", 13, True, TEAL),
        "• Extraction -> 14-Day Scorer -> Dietary Reconciliation Subagent.",
        ("Phase 3: Real-Time Output Stage", 13, True, TEAL),
        "• Dashboard Urgency Ring updates + Real-time Patient notification card."
    ]
)
img2_path = r"C:\Users\Shruti\.gemini\antigravity-ide\brain\4e55b1e4-91df-421f-8cc1-0a0d55313f34\ayuguard_user_workflow_diagram_1785176930297.png"
if os.path.exists(img2_path):
    s7.shapes.add_picture(img2_path, Inches(5.0), Inches(1.6), width=Inches(4.5))

# SLIDE 8
update_slide(
    prs.slides[7],
    "Live Production Deployments & Code Repository",
    [
        ("🌐 Live SPA Dashboard: https://ayuguard-ui-118454190848.asia-south1.run.app", 13, True, TEAL),
        ("🤖 Live Agent Orchestrator API: https://ayuguard-agent-118454190848.asia-south1.run.app", 13, True, TEAL),
        ("💻 GitHub Repository: https://github.com/shruti-patel06/AyuGuard", 13, True, TEAL),
        ("📍 Cloud Infrastructure: Google Cloud Run (asia-south1) | GCP Project: silken-dogfish-484814-g9", 12, False, GRAY)
    ]
)

# SLIDE 9
update_slide(
    prs.slides[8],
    "Technologies Used & System Design Justification",
    [
        "• Google ADK (Agent Development Kit): Modular Python framework establishing strict subagent boundaries and reliable tool execution loops.",
        "• Vertex AI (Gemini 2.5 Flash Lite): Enterprise LLM powering NLP symptom extraction, dietary reconciliation reasoning, and multimodal document analysis in us-central1.",
        "• Google Cloud Run: Managed serverless container hosting providing auto-scaling from zero to high concurrency in asia-south1.",
        "• Google Cloud Storage & Firestore: Encrypted PDF document vault and real-time notification feed sync.",
        "• Google Secret Manager: Secure API key & ADC credential management with enterprise compliance."
    ]
)

# SLIDE 10
update_slide(
    prs.slides[9],
    "Snapshots of Prototype: Caregiver Chat & Subagent Collaboration",
    [
        "• Conversational Logging: Accepts English, Hindi, and Hinglish input (e.g. 'Papa ko loose motion ho gaya hai').",
        "• Subagent Collaboration: dietary_reconciliation_agent automatically detects dietary conflicts and updates the care plan.",
        "• Real-Time Notification Bell: Flashes instantly when care plan revisions or critical alerts occur."
    ]
)

# SLIDE 11
update_slide(
    prs.slides[10],
    "Snapshots of Prototype: Health Dashboard & Medical Records Vault",
    [
        "• Glassmorphism Health Dashboard: Visual composite Urgency Ring, 14-day chronological timeline, and clinical precautions.",
        "• Multimodal Gemini Records Vault: Automatic extraction of abnormal lab parameters (HbA1c, blood glucose) from uploaded PDF reports.",
        "• Patient Chat View: Calm patient interface with real-time caregiver update cards."
    ]
)

# SLIDE 12
update_slide(
    prs.slides[11],
    "Conclusion & Future Roadmap",
    [
        ("Clinical Safety Boundary", 15, True, TEAL),
        "• AyuGuard supplements — never replaces — professional medical advice ('The math decides what is dangerous — the LLM only explains it').",
        ("Future Roadmap", 15, True, TEAL),
        "• Integration with smart IoT pill dispensers, voice-call IVR integration for non-smartphone seniors, and multi-caregiver family sync.",
        ("AyuGuard — Ambient Care, Powered by AI.", 16, True, NAVY)
    ]
)

prs.save('AyuGuard_Final_Submission_Deck.pptx')
print('AYUGUARD FINAL SUBMISSION DECK CREATED SUCCESSFULLY!')
